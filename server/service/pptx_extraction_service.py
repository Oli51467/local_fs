"""Utilities for extracting text and images from PPTX presentations."""

from __future__ import annotations

import hashlib
import tempfile
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Dict, List, Optional

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

@dataclass
class PptxExtractionResult:
    """Container for PPTX extraction outputs."""

    text: str
    images: List[Dict]
    temp_dir: Optional[Path]


IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.gif', '.tiff', '.bmp'}


class PptxExtractionError(Exception):
    """Raised when PPTX parsing fails."""


def _collect_paragraph_lines(shape) -> List[str]:
    lines: List[str] = []
    if not getattr(shape, "has_text_frame", False):
        return lines

    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        text = paragraph.text.strip()
        if text:
            lines.append(text)
    return lines


def _collect_table_lines(shape) -> List[str]:
    lines: List[str] = []
    if not getattr(shape, "has_table", False):
        return lines

    table = shape.table
    for row in table.rows:
        row_cells = []
        for cell in row.cells:
            cell_text = cell.text.strip()
            if cell_text:
                row_cells.append(cell_text)
        if row_cells:
            lines.append(" | ".join(row_cells))
    return lines


def _get_slide_title(slide) -> Optional[str]:
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None:
        title_lines = _collect_paragraph_lines(title_shape)
        if title_lines:
            return " ".join(title_lines).strip()

    for placeholder in getattr(slide, "placeholders", []):
        try:
            if placeholder.placeholder_format.type in {
                PP_PLACEHOLDER.TITLE,
                PP_PLACEHOLDER.CENTER_TITLE,
            }:
                text = placeholder.text.strip()
                if text:
                    return text
        except (AttributeError, KeyError):
            continue

    return None


def _get_notes_text(slide) -> Optional[str]:
    if not getattr(slide, "has_notes_slide", False):
        return None

    try:
        notes_slide = slide.notes_slide
    except KeyError:
        return None

    if not notes_slide or not notes_slide.notes_text_frame:
        return None

    text = notes_slide.notes_text_frame.text
    return text.strip() if text else None


def _shape_is_title(shape, title_shape) -> bool:
    return title_shape is not None and shape.shape_id == title_shape.shape_id


def _iter_shape_text(slide, title_shape) -> List[str]:
    collected: List[str] = []
    for shape in slide.shapes:
        if _shape_is_title(shape, title_shape):
            continue

        if getattr(shape, "has_text_frame", False):
            collected.extend(_collect_paragraph_lines(shape))
        elif getattr(shape, "has_table", False):
            collected.extend(_collect_table_lines(shape))
    return collected


def _extract_slide_text(slide, slide_index: int) -> Optional[str]:
    title_shape = getattr(slide.shapes, "title", None)
    title_text = _get_slide_title(slide)
    body_lines = _iter_shape_text(slide, title_shape)
    notes_text = _get_notes_text(slide)

    slide_lines: List[str] = []

    header_parts = [f"第{slide_index + 1}页"]
    if title_text:
        header_parts.append(title_text)
    slide_lines.append(" - ".join(header_parts))

    for line in body_lines:
        slide_lines.append(f"• {line}")

    if notes_text:
        slide_lines.append(f"备注：{notes_text}")

    meaningful_lines = [line for line in slide_lines if line.strip()]
    if not meaningful_lines:
        return None

    return "\n".join(meaningful_lines)


def _resolve_alt_text(shape) -> str:
    candidates = [
        getattr(shape, "alternative_text", None),
        getattr(shape, "alt_text", None),
        getattr(shape, "name", None),
    ]
    for candidate in candidates:
        if candidate and str(candidate).strip():
            return str(candidate).strip()
    return ""


def extract_pptx_text_and_images(pptx_path: Path) -> PptxExtractionResult:
    if not pptx_path.exists() or not pptx_path.is_file():
        raise PptxExtractionError(f"PPTX文件不存在: {pptx_path}")

    try:
        presentation = Presentation(str(pptx_path))
    except Exception as exc:  # pylint: disable=broad-except
        raise PptxExtractionError(f"无法读取PPTX文件: {exc}") from exc

    slide_texts: List[str] = []
    images: List[Dict] = []
    temp_dir: Optional[Path] = None
    seen_hashes: Dict[str, Path] = {}

    for slide_index, slide in enumerate(presentation.slides):
        slide_text = _extract_slide_text(slide, slide_index)
        if slide_text:
            slide_texts.append(slide_text)

        for shape_index, shape in enumerate(slide.shapes, start=1):
            image_descriptor = getattr(shape, "image", None)
            if image_descriptor is None:
                continue

            image_ext = f".{image_descriptor.ext.lower()}" if image_descriptor.ext else ""
            if image_ext not in IMAGE_EXTENSIONS:
                continue

            if temp_dir is None:
                temp_dir = Path(tempfile.mkdtemp(prefix="pptx_img_"))

            image_bytes = image_descriptor.blob
            fingerprint = hashlib.md5(image_bytes).hexdigest()
            if fingerprint in seen_hashes:
                continue

            dest_name = f"slide{slide_index + 1:03d}_shape{shape_index:03d}{image_ext}"
            dest_path = temp_dir / dest_name
            dest_path.write_bytes(image_bytes)

            width = height = None
            try:
                with Image.open(BytesIO(image_bytes)) as pil_image:
                    width, height = pil_image.size
            except Exception:  # pylint: disable=broad-except
                width = height = None

            stat_info = dest_path.stat()
            images.append(
                {
                    "source_path": dest_path,
                    "source_path_relative": None,
                    "line_number": slide_index + 1,
                    "chunk_index": slide_index,
                    "alt_text": _resolve_alt_text(shape),
                    "image_format": image_ext.lstrip('.'),
                    "image_size": stat_info.st_size,
                    "width": width,
                    "height": height,
                    "slide_index": slide_index,
                    "shape_id": getattr(shape, "shape_id", None),
                    "fingerprint": fingerprint,
                    "temp_dir": temp_dir,
                }
            )
            seen_hashes[fingerprint] = dest_path

    aggregated_text = "\n\n".join(slide_texts).strip()
    return PptxExtractionResult(text=aggregated_text, images=images, temp_dir=temp_dir)


__all__ = [
    "PptxExtractionResult",
    "PptxExtractionError",
    "extract_pptx_text_and_images",
]
