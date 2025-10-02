from pptx import Presentation
from pathlib import Path
import hashlib


def extract_text_from_pptx(file_path: str):
    prs = Presentation(file_path)
    all_text = []

    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            all_text.append(f"--- 第 {i+1} 页 ---\n" + "\n".join(slide_text))

    return "\n\n".join(all_text)


def extract_images_from_pptx(file_path: str, output_dir: str = "images"):
    prs = Presentation(file_path)
    output_dir = Path(output_dir)
    output_dir.mkdir(exist_ok=True)

    seen_hashes = {}  # 存放 {hash值: 文件名}
    image_count = 0

    for slide_id, slide in enumerate(prs.slides, start=1):
        for shape_id, shape in enumerate(slide.shapes, start=1):
            if hasattr(shape, "image"):
                image = shape.image
                image_bytes = image.blob

                # 计算哈希值（MD5）
                md5 = hashlib.md5(image_bytes).hexdigest()

                # 如果该图片已保存过，则跳过
                if md5 in seen_hashes:
                    print(f"⚙️ 发现重复图片（跳过）: {seen_hashes[md5]}")
                    continue

                # 否则保存
                image_ext = image.ext
                image_filename = f"slide{slide_id}_img{shape_id}.{image_ext}"
                image_path = output_dir / image_filename

                with open(image_path, "wb") as f:
                    f.write(image_bytes)

                seen_hashes[md5] = image_filename
                image_count += 1

    print(f"✅ 提取完成：共 {image_count} 张唯一图片，保存至 {output_dir.resolve()}")
    return seen_hashes


if __name__ == "__main__":
    extract_images_from_pptx(
        "../../data/文档/副本DeepSeek部署方案_20250407.pptx",
        output_dir="../../data/images",
    )
    text = extract_text_from_pptx("../../data/文档/副本DeepSeek部署方案_20250407.pptx")
    print(text)
